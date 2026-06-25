import os
import json
import argparse
import re
from datetime import datetime
from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def sanitize_filename(filename: str) -> str:
    """
    파일명으로 사용할 수 없는 특수문자를 제거하거나 언더스코어로 변경합니다.
    """
    return re.sub(r'[\/\\\:\*\?\"\<\>\|]', '_', filename)

def create_presentation_for_paper(summary_json_path: str, base_dir: str):
    """
    단일 논문 요약 JSON을 바탕으로 해당 논문 전용 PPTX 파일을 빌드합니다.
    """
    if not os.path.exists(summary_json_path):
        print(f"[PPT_GEN] Error: Summary file not found at {summary_json_path}")
        return

    # JSON 데이터 로드
    try:
        with open(summary_json_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        print(f"[PPT_GEN] Failed to read JSON file: {e}")
        return

    original_title = data.get("original_title", "Untitled Paper")
    korean_title = data.get("korean_title", "제목 없음")
    sections = data.get("sections", [])
    equations = data.get("equations", [])

    # 파일 경로 매핑
    # 요약 JSON 파일과 동일한 경로 기준
    json_dir = os.path.dirname(summary_json_path)
    base_name = os.path.basename(summary_json_path).replace("_summary.json", "")
    
    # SDXL 일러스트 및 사용자 지정 이미지 경로 (paper_illustration 폴더 등 구조에 맞게 탐색)
    illustration_dir = json_dir.replace("paper_summary", "paper_illustration")
    illustration_path = os.path.join(illustration_dir, f"{base_name}_illustration.png")
    
    # 사용자가 직접 캡처하여 넣은 논문 핵심 구조 이미지 (존재할 경우 최우선 사용)
    custom_core_image = os.path.join(json_dir, f"{base_name}_core.png")
    
    # PDF 추출 이미지 폴더
    pdf_images_dir = os.path.join(json_dir, f"{base_name}_images")
    pdf_images = []
    if os.path.exists(pdf_images_dir):
        # 추출된 이미지들 정렬
        pdf_images = [os.path.join(pdf_images_dir, f) for f in os.listdir(pdf_images_dir) 
                      if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        # 페이지 순으로 정렬하기 위해 정렬 시도
        pdf_images.sort()

    print(f"[PPT_GEN] Generating PPT for: {korean_title}")
    print(f" -> Found {len(pdf_images)} extracted PDF images.")

    # 새로운 PPT 문서 생성
    prs = Presentation()
    
    # 16:9 슬라이드 크기 설정 (표준 와이드스크린)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 제목 슬라이드 추가 (Layout 6: Blank)
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 제목 텍스트 박스 커스텀 배치 (글자 길이에 맞춘 레이아웃 안정성 확보)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # 한글 제목 (대제목)
    p_title = tf.paragraphs[0]
    p_title.text = korean_title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(10, 30, 80)
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(20)

    # 영어 원제 (소제목)
    p_sub = tf.add_paragraph()
    p_sub.text = f"Original Title: {original_title}"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = RGBColor(120, 120, 120)
    p_sub.alignment = PP_ALIGN.LEFT
    p_sub.space_after = Pt(30)
    
    # 메타 정보 (날짜 등)
    p_meta = tf.add_paragraph()
    p_meta.text = f"발표 및 요약 일자: {datetime.now().strftime('%Y-%m-%d')} | AI 논문 분석 시스템"
    p_meta.font.name = "Arial"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = RGBColor(150, 150, 150)
    p_meta.alignment = PP_ALIGN.LEFT

    # 2. 본문 요약 슬라이드 추가 (섹션별 루프)
    blank_layout = prs.slide_layouts[6]

    used_images = set()

    for idx, sec in enumerate(sections):
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes
        
        sec_title = sec.get("title", f"섹션 {idx+1}")
        sec_content = sec.get("content", "")
        
        # 슬라이드 대제목 추가
        title_box = shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf_title = title_box.text_frame
        p_sec_title = tf_title.paragraphs[0]
        p_sec_title.text = sec_title
        p_sec_title.font.name = "Arial"
        p_sec_title.font.size = Pt(24)
        p_sec_title.font.bold = True
        p_sec_title.font.color.rgb = RGBColor(0, 102, 204)
        
        # 이미지 매핑 후보 찾기
        img_to_place = None
        
        # 1. 사용자가 직접 추가한 핵심 구조 이미지가 존재하면 가장 중요한 제안 방법론(idx=1)에 배치
        if idx == 1 and os.path.exists(custom_core_image):
            img_to_place = custom_core_image
        # 2. PDF에서 추출된 이미지가 있을 경우
        elif pdf_images:
            # 서론(idx=0) -> 첫 번째 이미지
            # 방법론(idx=1) -> 중간 이미지
            # 결과/결론(idx=2) -> 마지막 이미지
            target_img_idx = 0
            if idx == 1 and len(pdf_images) > 1:
                target_img_idx = len(pdf_images) // 2
            elif idx == 2 and len(pdf_images) > 2:
                target_img_idx = len(pdf_images) - 1
                
            candidate = pdf_images[target_img_idx]
            if os.path.exists(candidate):
                img_to_place = candidate

        # 3. PDF 이미지가 매핑되지 않은 경우, SDXL 일러스트 이미지를 대체 적용
        if not img_to_place and os.path.exists(illustration_path):
            img_to_place = illustration_path
            
        # 좌측 텍스트 영역 너비 결정 (이미지가 있으면 6.0인치, 없으면 11.833인치로 확장)
        text_width = Inches(6.0) if img_to_place else Inches(11.833)
        
        # 텍스트 영역 (본문 요약)
        content_box = shapes.add_textbox(Inches(0.75), Inches(1.5), text_width, Inches(5.0))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        p_content = tf_content.paragraphs[0]
        p_content.text = sec_content
        p_content.font.name = "Arial"
        p_content.font.size = Pt(15)
        p_content.font.color.rgb = RGBColor(50, 50, 50)
        p_content.line_spacing = 1.3
        
        # 이미지 배치
        if img_to_place:
            try:
                shapes.add_picture(img_to_place, Inches(7.5), Inches(1.5), width=Inches(5.0))
                used_images.add(img_to_place)
                print(f"[PPT_GEN] Placed image on Slide '{sec_title}': {os.path.basename(img_to_place)}")
            except Exception as e:
                print(f"[PPT_GEN] Warning: Failed to add image: {e}")

    # 3. 수식 전용 슬라이드 추가 (수식이 존재하는 경우에만 필수 추가)
    if equations:
        out_dir_tmp = os.path.join(illustration_dir, "equations")
        os.makedirs(out_dir_tmp, exist_ok=True)
        
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes
        
        # 슬라이드 대제목 추가
        title_box = shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf_title = title_box.text_frame
        p_sec_title = tf_title.paragraphs[0]
        p_sec_title.text = "핵심 수식 및 이론적 근거 (Mathematical Formulation)"
        p_sec_title.font.name = "Arial"
        p_sec_title.font.size = Pt(24)
        p_sec_title.font.bold = True
        p_sec_title.font.color.rgb = RGBColor(120, 30, 80)
        
        current_y = 1.5
        
        for eq_idx, eq_data in enumerate(equations):
            formula = eq_data.get("formula", "")
            explanation = eq_data.get("explanation", "")
            
            if not formula:
                continue
                
            if current_y > 5.5:  # 슬라이드를 벗어날 경우 새 슬라이드 추가
                slide = prs.slides.add_slide(blank_layout)
                shapes = slide.shapes
                title_box = shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
                p_sec_title = title_box.text_frame.paragraphs[0]
                p_sec_title.text = "핵심 수식 및 이론적 근거 (계속)"
                p_sec_title.font.name = "Arial"
                p_sec_title.font.size = Pt(24)
                p_sec_title.font.bold = True
                p_sec_title.font.color.rgb = RGBColor(120, 30, 80)
                current_y = 1.5
                
            # 수식 라벨
            lbl_box = shapes.add_textbox(Inches(0.75), Inches(current_y), Inches(11.833), Inches(0.4))
            p_lbl = lbl_box.text_frame.paragraphs[0]
            p_lbl.text = f"Formula {eq_idx + 1}:"
            p_lbl.font.name = "Arial"
            p_lbl.font.size = Pt(14)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = RGBColor(0, 102, 204)
            
            current_y += 0.4
            
            # 수식 이미지 생성 및 삽입
            try:
                clean_formula = formula.strip().strip('$')
                fig, ax = plt.subplots(figsize=(10, 1.5))
                ax.text(0.01, 0.5, f"${clean_formula}$", fontsize=22, va='center')
                ax.axis('off')
                eq_img_path = os.path.join(out_dir_tmp, f"temp_eq_{eq_idx}.png")
                fig.savefig(eq_img_path, bbox_inches='tight', transparent=True, pad_inches=0.1)
                plt.close(fig)
                
                pic = shapes.add_picture(eq_img_path, Inches(1.0), Inches(current_y))
                if pic.width > Inches(10.0):
                    pic.width = Inches(10.0)
                current_y += (pic.height.inches + 0.1)
            except Exception as e:
                print(f"[PPT_GEN] Warning: Formula rendering failed: {e}")
                # 렌더링 실패 시 텍스트 렌더링
                eq_box = shapes.add_textbox(Inches(1.0), Inches(current_y), Inches(10.0), Inches(0.5))
                p_form = eq_box.text_frame.paragraphs[0]
                p_form.text = f"   {formula}   "
                p_form.font.name = "Courier New"
                p_form.font.size = Pt(20)
                current_y += 0.6
            
            # 설명
            exp_box = shapes.add_textbox(Inches(0.75), Inches(current_y), Inches(11.833), Inches(0.6))
            exp_box.text_frame.word_wrap = True
            p_exp = exp_box.text_frame.paragraphs[0]
            p_exp.text = f"💡 해석: {explanation}"
            p_exp.font.name = "Arial"
            p_exp.font.size = Pt(13)
            p_exp.font.color.rgb = RGBColor(80, 80, 80)
            
            current_y += 0.7

    # 3.5 추가 시각 자료 슬라이드 (사용되지 않은 PDF 이미지들)
    remaining_images = [img for img in pdf_images if img not in used_images]
    if remaining_images:
        for img_idx, img_path in enumerate(remaining_images):
            slide = prs.slides.add_slide(blank_layout)
            shapes = slide.shapes
            
            # Title
            title_box = shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
            p_title = title_box.text_frame.paragraphs[0]
            p_title.text = f"추가 시각 자료 (Figure {img_idx + 1})"
            p_title.font.name = "Arial"
            p_title.font.size = Pt(24)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(0, 102, 204)
            
            # Image
            try:
                pic = shapes.add_picture(img_path, Inches(1.0), Inches(1.5))
                max_width = Inches(11.333)
                max_height = Inches(5.5)
                
                if pic.width > max_width or pic.height > max_height:
                    width_ratio = max_width / pic.width
                    height_ratio = max_height / pic.height
                    scale = min(width_ratio, height_ratio)
                    pic.width = int(pic.width * scale)
                    pic.height = int(pic.height * scale)
                
                # Center the image
                pic.left = int((prs.slide_width - pic.width) / 2)
                pic.top = int(Inches(1.5) + (max_height - pic.height) / 2)
            except Exception as e:
                print(f"[PPT_GEN] Warning: Failed to add remaining image: {e}")
            
    # 4. 저장 폴더 생성 및 PPTX 저장
    out_dir = os.path.join(base_dir, "output")
    os.makedirs(out_dir, exist_ok=True)
    
    safe_title = sanitize_filename(korean_title)
    out_file = os.path.join(out_dir, f"{safe_title}.pptx")
    
    prs.save(out_file)
    print(f"[PPT_GEN] Success! Single paper presentation saved to: {out_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="단일 논문 요약 JSON을 받아 개별 PPTX 빌드")
    parser.add_argument("--json", type=str, required=True, help="빌드할 단일 요약 JSON 파일 경로")
    
    args = parser.parse_args()
    root_directory = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    
    create_presentation_for_paper(os.path.abspath(args.json), root_directory)
